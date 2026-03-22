"""
ppt_generator.py
- One slide per section (no splitting)
- All section content on that slide
- Font size 24 for bullet content only
- Template visuals preserved exactly
"""

import copy, re, shutil, zipfile
import lxml.etree as etree
from pptx import Presentation
from pptx.oxml.ns import qn


# ─────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────

def generate_ppt(doc_data: dict, template_data: dict,
                 template_path: str, output_path: str) -> str:

    title    = doc_data.get('title', 'Presentation')
    sections = [s for s in doc_data.get('sections', [])
                if s.get('heading') or s.get('points')]

    prs = Presentation(template_path)
    sldIdLst = list(prs.slides._sldIdLst)

    # Get the two template slide parts
    title_part   = prs.slides.part.related_part(sldIdLst[0].get(qn('r:id')))
    content_part = prs.slides.part.related_part(sldIdLst[1].get(qn('r:id')))

    # ── Build modified XMLs ──────────────────────────────
    # 1. Title slide
    title_xml = copy.deepcopy(title_part._element)
    _set_title_slide_text(title_xml, title)

    # 2. One content slide per section — ALL points on one slide, no splitting
    content_xmls = []
    for section in sections:
        heading = section.get('heading', '')
        points  = section.get('points', [])
        xml = copy.deepcopy(content_part._element)
        _set_content_slide_text(xml, title, heading, points)
        content_xmls.append(xml)

    # ── Write to file ────────────────────────────────────
    _write_pptx(template_path, output_path, title_xml, content_xmls)

    total = 1 + len(content_xmls)
    print(f"  Saved: {output_path} ({total} slides)")
    return output_path


# ─────────────────────────────────────────────────────────
# Text Setters
# ─────────────────────────────────────────────────────────

def _set_title_slide_text(xml, title: str):
    """Find the main chapter title shape and replace its text."""
    for sp in xml.iter():
        if sp.tag not in (qn('p:sp'), qn('p:grpSp')):
            continue
        all_text = ''.join(
            t.text for t in sp.findall('.//' + qn('a:t')) if t.text
        )
        if 'chapter' in all_text.lower() or 'robot' in all_text.lower() or 'start' in all_text.lower():
            _replace_multirun_text(sp, title)
            return


def _replace_multirun_text(sp, new_text: str):
    """Replace a shape that has text split across multiple runs."""
    txBody = sp.find('.//' + qn('p:txBody'))
    if txBody is None:
        return

    # Grab formatting from first run
    first_rPr = None
    for rPr in txBody.findall('.//' + qn('a:rPr')):
        first_rPr = copy.deepcopy(rPr)
        break

    # Remove all paragraphs
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    # Single paragraph, single run
    p = etree.SubElement(txBody, qn('a:p'))
    r = etree.SubElement(p, qn('a:r'))
    if first_rPr is not None:
        r.append(first_rPr)
    t = etree.SubElement(r, qn('a:t'))
    t.text = new_text


def _set_content_slide_text(xml, doc_title: str, heading: str, points: list):
    """
    Content slide shape positions (Chapter_8.pptx):
      header bar   : top < 700k,   left > 1M,  w > 10M  → doc title
      section label: top 500k–1.5M, left > 3M, w > 5M   → section name
      content body : top > 1.5M,   top < 9M,   w > 10M  → bullets
    """
    spTree = xml.find('.//' + qn('p:spTree'))
    for sp in spTree.findall(qn('p:sp')):
        top, left, w, h = _get_pos(sp)

        # Header bar — replace with doc title (keep template formatting)
        if top < 700000 and left > 1000000 and w > 10000000:
            _replace_all_text_in_sp(sp, doc_title)

        # Section label in teal bar
        elif 500000 < top < 1500000 and left > 3000000 and w > 5000000:
            clean = re.sub(r'^\d+[\.\)]\s*', '', heading).strip()
            _replace_all_text_in_sp(sp, clean)

        # Content body — rebuild with heading + all bullets at font size 24
        elif top > 1500000 and top < 9000000 and w > 10000000:
            _set_content_body(sp, heading, points, font_size=24)


def _replace_all_text_in_sp(sp, new_text: str):
    """Replace all text in a shape keeping first run's formatting."""
    txBody = sp.find(qn('p:txBody'))
    if txBody is None:
        return

    first_rPr = None
    for rPr in txBody.findall('.//' + qn('a:rPr')):
        first_rPr = copy.deepcopy(rPr)
        break

    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    p = etree.SubElement(txBody, qn('a:p'))
    r = etree.SubElement(p, qn('a:r'))
    if first_rPr is not None:
        r.append(first_rPr)
    t = etree.SubElement(r, qn('a:t'))
    t.text = new_text


def _set_content_body(sp, heading: str, points: list, font_size: int = 24):
    """
    Build content body:
    - Bold heading line (font_size)
    - Blank line
    - All bullet points (font_size, teal bullet)
    ALL content uses font_size=24 as specified.
    """
    txBody = sp.find(qn('p:txBody'))
    if txBody is None:
        return

    # Get font name from existing content
    font_name = 'Calibri'
    for rPr in txBody.findall('.//' + qn('a:rPr')):
        latin = rPr.find(qn('a:latin'))
        if latin is not None:
            font_name = latin.get('typeface', 'Calibri')
            break

    # Remove all existing paragraphs
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    def mk_color(hex_str):
        el = etree.Element(qn('a:solidFill'))
        etree.SubElement(el, qn('a:srgbClr'), attrib={'val': hex_str.upper()})
        return el

    def mk_latin(name):
        return etree.Element(qn('a:latin'), attrib={'typeface': name})

    def add_para(text, bold=False, color='1A1A2E', bullet=False):
        p = etree.SubElement(txBody, qn('a:p'))

        pPr = etree.SubElement(p, qn('a:pPr'))
        pPr.set('algn', 'l')
        if bullet:
            pPr.set('marL', '342900')
            pPr.set('indent', '-342900')
            buClr = etree.SubElement(pPr, qn('a:buClr'))
            buClr.append(mk_color('1ABEBC'))
            etree.SubElement(pPr, qn('a:buFont'), attrib={'typeface': font_name})
            etree.SubElement(pPr, qn('a:buChar'), attrib={'char': '•'})
        else:
            etree.SubElement(pPr, qn('a:buNone'))

        if text:
            r = etree.SubElement(p, qn('a:r'))
            rPr = etree.SubElement(r, qn('a:rPr'), attrib={
                'lang': 'en-US',
                'dirty': '0',
                'b': '1' if bold else '0',
                'sz': str(font_size * 100)   # Always 24pt for ALL content
            })
            rPr.append(mk_color(color))
            rPr.append(mk_latin(font_name))
            t = etree.SubElement(r, qn('a:t'))
            t.text = text

    # Section heading (bold, navy)
    add_para(heading, bold=True, color='113264', bullet=False)

    # Blank separator line
    etree.SubElement(txBody, qn('a:p'))

    # All bullet points — every single one, no limit
    for pt in points:
        if pt and pt.strip():
            add_para(pt.strip(), bold=False, color='2C3E50', bullet=True)


# ─────────────────────────────────────────────────────────
# Shape helpers
# ─────────────────────────────────────────────────────────

def _get_pos(sp):
    xfrm = sp.find('.//' + qn('a:xfrm'))
    if xfrm is None:
        return 0, 0, 0, 0
    off = xfrm.find(qn('a:off'))
    ext = xfrm.find(qn('a:ext'))
    top  = int(off.get('y', 0)) if off is not None else 0
    left = int(off.get('x', 0)) if off is not None else 0
    w    = int(ext.get('cx', 0)) if ext is not None else 0
    h    = int(ext.get('cy', 0)) if ext is not None else 0
    return top, left, w, h


# ─────────────────────────────────────────────────────────
# PPTX File Writer (direct zip manipulation)
# ─────────────────────────────────────────────────────────

def _write_pptx(template_path: str, output_path: str,
                title_xml, content_xmls: list):

    import re as re2

    shutil.copy2(template_path, output_path)

    with zipfile.ZipFile(output_path, 'r') as zin:
        names = zin.namelist()
        files = {}
        for name in names:
            files[name] = zin.read(name)

    slide_files = sorted(
        [n for n in names if re2.match(r'ppt/slides/slide\d+\.xml$', n)],
        key=lambda x: int(re2.search(r'\d+', x.split('/')[-1]).group())
    )
    rel_files = sorted(
        [n for n in names if re2.match(r'ppt/slides/_rels/slide\d+\.xml\.rels$', n)],
        key=lambda x: int(re2.search(r'\d+', x.split('/')[-1]).group())
    )

    template_slide1_rel = files[rel_files[0]]
    template_slide2_rel = files[rel_files[1]]

    all_slides_xml = [
        etree.tostring(title_xml, xml_declaration=True,
                       encoding='UTF-8', standalone=True)
    ] + [
        etree.tostring(cx, xml_declaration=True,
                       encoding='UTF-8', standalone=True)
        for cx in content_xmls
    ]

    new_files = dict(files)

    for sf in slide_files:
        del new_files[sf]
    for rf in rel_files:
        del new_files[rf]

    for i, xml_bytes in enumerate(all_slides_xml, start=1):
        new_files[f'ppt/slides/slide{i}.xml'] = xml_bytes
        new_files[f'ppt/slides/_rels/slide{i}.xml.rels'] = (
            template_slide1_rel if i == 1 else template_slide2_rel
        )

    # Update presentation.xml
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    ns_rel = 'http://schemas.openxmlformats.org/package/2006/relationships'
    RT_SLIDE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'

    prs_xml = etree.fromstring(files['ppt/presentation.xml'])
    sldIdLst = prs_xml.find(f'{{{ns_p}}}sldIdLst')
    for child in list(sldIdLst):
        sldIdLst.remove(child)

    prs_rels_xml = etree.fromstring(files['ppt/_rels/presentation.xml.rels'])

    # Get existing rId map for slides
    existing_rids = {}
    for rel in prs_rels_xml.findall(f'{{{ns_rel}}}Relationship'):
        target = rel.get('Target', '')
        rid    = rel.get('Id', '')
        if 'slides/slide' in target:
            m = re2.search(r'slide(\d+)\.xml', target)
            if m:
                existing_rids[int(m.group(1))] = rid

    all_rid_nums = []
    for rel in prs_rels_xml.findall(f'{{{ns_rel}}}Relationship'):
        m = re2.search(r'(\d+)', rel.get('Id', ''))
        if m:
            all_rid_nums.append(int(m.group(1)))
    max_rid_num = max(all_rid_nums) if all_rid_nums else 10

    slide_rid_map = {}
    for i in range(1, len(all_slides_xml) + 1):
        if i in existing_rids:
            slide_rid_map[i] = existing_rids[i]
        else:
            max_rid_num += 1
            slide_rid_map[i] = f'rId{max_rid_num}'

    max_slide_id = 255
    for el in prs_xml.iter():
        sid = el.get('id')
        if sid and el.tag.endswith('sldId'):
            try:
                max_slide_id = max(max_slide_id, int(sid))
            except ValueError:
                pass

    for i in range(1, len(all_slides_xml) + 1):
        max_slide_id += 1
        etree.SubElement(sldIdLst, f'{{{ns_p}}}sldId', attrib={
            'id': str(max_slide_id),
            f'{{{ns_r}}}id': slide_rid_map[i]
        })

    new_files['ppt/presentation.xml'] = etree.tostring(
        prs_xml, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Update presentation.xml.rels
    for rel in list(prs_rels_xml):
        if rel.get('Type') == RT_SLIDE:
            prs_rels_xml.remove(rel)

    for i in range(1, len(all_slides_xml) + 1):
        etree.SubElement(prs_rels_xml, f'{{{ns_rel}}}Relationship', attrib={
            'Id': slide_rid_map[i],
            'Type': RT_SLIDE,
            'Target': f'slides/slide{i}.xml'
        })

    new_files['ppt/_rels/presentation.xml.rels'] = etree.tostring(
        prs_rels_xml, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Update [Content_Types].xml
    ct_xml = etree.fromstring(files['[Content_Types].xml'])
    CT_NS = 'http://schemas.openxmlformats.org/package/2006/content-types'
    SLIDE_CT = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'

    for override in list(ct_xml):
        pn = override.get('PartName', '')
        if re2.match(r'/ppt/slides/slide\d+\.xml', pn):
            ct_xml.remove(override)

    for i in range(1, len(all_slides_xml) + 1):
        etree.SubElement(ct_xml, f'{{{CT_NS}}}Override', attrib={
            'PartName': f'/ppt/slides/slide{i}.xml',
            'ContentType': SLIDE_CT
        })

    new_files['[Content_Types].xml'] = etree.tostring(
        ct_xml, xml_declaration=True, encoding='UTF-8', standalone=True)

    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in new_files.items():
            zout.writestr(name, data)
