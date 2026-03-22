"""
template_reader.py — Analyze an uploaded .pptx template
Extracts colors, fonts, layouts, placeholder positions
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


def analyze_template(pptx_path: str) -> dict:
    """
    Analyze a .pptx file and return its design characteristics.
    """
    prs = Presentation(pptx_path)

    result = {
        "slide_count":        len(prs.slides),
        "width":              prs.slide_width,
        "height":             prs.slide_height,
        "layouts":            [],
        "colors":             _extract_colors(prs),
        "fonts":              _extract_fonts(prs),
        "title_slide_idx":    0,
        "content_slide_idx":  1,
        "slide_designs":      [],
    }

    # Analyze each slide
    for i, slide in enumerate(prs.slides):
        design = _analyze_slide(slide, i)
        result["slide_designs"].append(design)

        if design["has_title"] and not design["has_body"] and i == 0:
            result["title_slide_idx"] = i
        if design["has_title"] and design["has_body"] and result["content_slide_idx"] == 1:
            result["content_slide_idx"] = i

    # Extract layout info safely
    for j, layout in enumerate(prs.slide_layouts):
        layout_info = {"index": j, "name": layout.name, "placeholders": []}
        try:
            for ph in layout.placeholders:
                layout_info["placeholders"].append({
                    "idx":  ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
        except Exception:
            pass
        result["layouts"].append(layout_info)

    print(f"  Template: {result['slide_count']} slides | "
          f"primary={result['colors']['primary']} | "
          f"font={result['fonts']['heading']}")
    return result


# ── Safe placeholder check ─────────────────────────────────
def _is_placeholder(shape) -> bool:
    """Safely check if a shape is a placeholder."""
    try:
        return (shape.is_placeholder and
                shape.placeholder_format is not None)
    except Exception:
        return False


def _get_placeholder_idx(shape) -> int:
    """Safely get placeholder index, returns -1 if not a placeholder."""
    try:
        if _is_placeholder(shape):
            return shape.placeholder_format.idx
    except Exception:
        pass
    return -1


# ── Slide Analyzer ─────────────────────────────────────────
def _analyze_slide(slide, idx: int) -> dict:
    design = {
        "index":            idx,
        "has_title":        False,
        "has_body":         False,
        "background_color": None,
        "title_position":   None,
        "body_position":    None,
        "text_boxes":       [],
    }

    # Background color
    try:
        fill = slide.background.fill
        if fill.type is not None:
            design["background_color"] = str(fill.fore_color.rgb)
    except Exception:
        pass

    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue

            tf        = shape.text_frame
            font_info = _get_font_info(tf)
            ph_idx    = _get_placeholder_idx(shape)
            name_l    = shape.name.lower()

            box = {
                "name":   shape.name,
                "left":   shape.left,
                "top":    shape.top,
                "width":  shape.width,
                "height": shape.height,
                "text":   tf.text.strip()[:80],
                "font":   font_info,
            }
            design["text_boxes"].append(box)

            # Detect title
            is_title = (ph_idx == 0 or "title" in name_l)
            # Detect body
            is_body  = (ph_idx in (1, 2) or
                        any(k in name_l for k in ("body", "content", "text")))

            if is_title:
                design["has_title"]       = True
                design["title_position"]  = {**box, "font": font_info}
            elif is_body:
                design["has_body"]        = True
                design["body_position"]   = {**box, "font": font_info}

        except Exception:
            continue  # skip any problematic shape

    return design


# ── Color Extractor ────────────────────────────────────────
def _extract_colors(prs: Presentation) -> dict:
    colors = {
        "primary":    "0A2D5E",
        "secondary":  "1ABEBC",
        "text":       "1A1A2E",
        "background": "FFFFFF",
        "accent":     "E74C3C",
    }

    all_colors = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # Shape fill
            try:
                fill = shape.fill
                if fill.type is not None:
                    c = str(fill.fore_color.rgb)
                    if c and len(c) == 6:
                        all_colors.append(c)
            except Exception:
                pass

            # Text runs
            if not shape.has_text_frame:
                continue
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.type is not None:
                                c = str(run.font.color.rgb)
                                if c and len(c) == 6:
                                    all_colors.append(c)
                        except Exception:
                            pass
            except Exception:
                pass

    # Pick most common non-trivial colors
    from collections import Counter
    skip = {"FFFFFF", "000000", "FEFFFF", "000001", "FFFFFE"}
    filtered = [c for c in all_colors if c.upper() not in skip]

    if filtered:
        top = Counter(filtered).most_common(5)
        if len(top) >= 1: colors["primary"]   = top[0][0]
        if len(top) >= 2: colors["secondary"] = top[1][0]
        if len(top) >= 3: colors["accent"]    = top[2][0]

    return colors


# ── Font Extractor ─────────────────────────────────────────
def _extract_fonts(prs: Presentation) -> dict:
    fonts = {
        "heading":       "Calibri",
        "body":          "Calibri",
        "heading_size":  24,
        "body_size":     14,
        "heading_bold":  True,
        "heading_color": "FFFFFF",
        "body_color":    "1A1A2E",
    }

    heading_fonts, body_fonts   = [], []
    heading_sizes, body_sizes   = [], []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                is_title = "title" in shape.name.lower() or _get_placeholder_idx(shape) == 0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            fname = run.font.name
                            fsize = run.font.size
                            if fname:
                                if is_title:
                                    heading_fonts.append(fname)
                                    if fsize:
                                        heading_sizes.append(int(fsize / 12700))
                                else:
                                    body_fonts.append(fname)
                                    if fsize:
                                        body_sizes.append(int(fsize / 12700))
                        except Exception:
                            pass
            except Exception:
                pass

    from collections import Counter
    if heading_fonts:
        fonts["heading"]      = Counter(heading_fonts).most_common(1)[0][0]
    if body_fonts:
        fonts["body"]         = Counter(body_fonts).most_common(1)[0][0]
    if heading_sizes:
        fonts["heading_size"] = max(16, Counter(heading_sizes).most_common(1)[0][0])
    if body_sizes:
        s = Counter(body_sizes).most_common(1)[0][0]
        fonts["body_size"]    = max(10, min(s, 20))

    return fonts


# ── Font Info Helper ───────────────────────────────────────
def _get_font_info(text_frame) -> dict:
    info = {"name": "Calibri", "size": 14, "bold": False, "color": "1A1A2E"}
    try:
        for para in text_frame.paragraphs:
            for run in para.runs:
                try:
                    f = run.font
                    if f.name:  info["name"] = f.name
                    if f.size:  info["size"] = int(f.size / 12700)
                    if f.bold is not None: info["bold"] = f.bold
                    try:
                        if f.color and f.color.type is not None:
                            info["color"] = str(f.color.rgb)
                    except Exception:
                        pass
                    return info
                except Exception:
                    pass
    except Exception:
        pass
    return info
